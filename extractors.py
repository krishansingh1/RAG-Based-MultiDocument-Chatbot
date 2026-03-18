# extractors.py
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import openpyxl


def extract_from_pdf(file):
    """Extract text from a PDF file."""
    reader = PdfReader(file)
    text = ""
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text + "\n"
    return text


def extract_from_docx(file):
    """Extract text from a Word .docx file."""
    doc = Document(file)
    text = ""
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            text += paragraph.text + "\n"
    # Also extract text from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells]
            text += " | ".join(row_text) + "\n"
    return text


def extract_from_pptx(file):
    """Extract text from a PowerPoint .pptx file."""
    prs = Presentation(file)
    text = ""
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = f"--- Slide {slide_num} ---\n"
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        slide_text += paragraph.text + "\n"
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    slide_text += " | ".join(row_text) + "\n"
        text += slide_text + "\n"
    return text


def extract_from_excel(file):
    """Extract text from an Excel .xlsx file."""
    wb = openpyxl.load_workbook(file)
    text = ""
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        text += f"--- Sheet: {sheet_name} ---\n"
        for row in sheet.iter_rows(values_only=True):
            row_text = [str(cell) if cell is not None else "" for cell in row]
            text += " | ".join(row_text) + "\n"
        text += "\n"
    return text


def extract_text(uploaded_file):
    """
    Route to the correct extractor based on file extension.
    This is the single entry point — the rest of the app
    just calls this function and doesn't care about file type.
    """
    name = uploaded_file.name.lower()

    if name.endswith(".pdf"):
        return extract_from_pdf(uploaded_file)
    elif name.endswith(".docx"):
        return extract_from_docx(uploaded_file)
    elif name.endswith(".pptx"):
        return extract_from_pptx(uploaded_file)
    elif name.endswith(".xlsx"):
        return extract_from_excel(uploaded_file)
    else:
        raise ValueError(f"Unsupported file type: {name}")