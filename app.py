import os

from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
from openpyxl import load_workbook


def get_pdf_text(file_path):
    text = ""
    pdf_reader = PdfReader(file_path)
    for page in pdf_reader.pages:
        text += page.extract_text() or ""
    return text


def get_pptx_text(file_path):
    text = ""
    prs = Presentation(file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text += paragraph.text + "\n"
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        text += cell.text + "\t"
                    text += "\n"
    return text


def get_docx_text(file_path):
    text = ""
    doc = Document(file_path)
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                text += cell.text + "\t"
            text += "\n"
    return text


def get_xlsx_text(file_path):
    text = ""
    wb = load_workbook(file_path, data_only=True)
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        text += f"--- Sheet: {sheet_name} ---\n"
        for row in sheet.iter_rows(values_only=True):
            row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
            text += row_text + "\n"
    return text


EXTRACTORS = {
    ".pdf": get_pdf_text,
    ".pptx": get_pptx_text,
    ".docx": get_docx_text,
    ".xlsx": get_xlsx_text,
}


def get_text_from_files(file_paths):
    all_text = ""
    for file_path in file_paths:
        ext = os.path.splitext(file_path)[1].lower()
        extractor = EXTRACTORS.get(ext)
        if extractor is None:
            print(f"Unsupported file type: {ext} ({file_path})")
            continue
        all_text += extractor(file_path) + "\n"
    return all_text


if __name__ == "__main__":
    test_files = ["Individual_Report_Solution_Internet_and_Cloud_Computing.pdf"]
    text = get_text_from_files(test_files)
    print(f"Extracted {len(text)} characters")
    